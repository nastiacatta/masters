#!/usr/bin/env python3
"""
Build the thesis defence PowerPoint from the Imperial College template.

Produces presentation/slides.pptx with 14 slides:
  PROBLEM:    1 Title, 2 Motivation, 3 Markets, 4 Existing Work, 5 Gap
  SOLUTION:   6 Mechanism, 7 Skill Signal, 8 Architecture
  VALIDATION: 9 Correctness, 10 Deposit Design, 11 Weight Rules,
              12 Skill Recovery, 13 Strategic Robustness, 14 Closing

Usage: python presentation/build_slides.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

TEMPLATE = "presentation/template.pptx"
OUTPUT = "presentation/slides.pptx"
PLOTS = "presentation/plots"

# Imperial colours (from template)
NAVY = RGBColor(0x00, 0x2A, 0x4E)
TEAL = RGBColor(0x00, 0x9C, 0xBC)
SLATE = RGBColor(0x5B, 0x6B, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF0, 0xF0, 0xF0)


def get_layout(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise ValueError(f"Layout '{name}' not found")


def set_text(placeholder, text, font_size=None, bold=False, color=None, alignment=None):
    """Set text on a placeholder, clearing existing content."""
    placeholder.text = ""
    tf = placeholder.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if font_size:
        p.font.size = Pt(font_size)
    if bold:
        p.font.bold = True
    if color:
        p.font.color.rgb = color
    if alignment:
        p.alignment = alignment


def add_bullet_slide(prs, layout_name, title, subtitle, bullets):
    """Add a slide with title, subtitle, and bullet points."""
    slide = prs.slides.add_slide(get_layout(prs, layout_name))
    # Set title
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, title, font_size=28, bold=True)
        elif ph.placeholder_format.idx == 13:
            set_text(ph, subtitle, font_size=14, color=SLATE)
        elif ph.placeholder_format.idx == 1:
            ph.text = ""
            tf = ph.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(16)
                p.font.color.rgb = SLATE
                p.space_after = Pt(6)
                p.level = 0
    return slide


def add_image_to_slide(slide, image_path, left, top, width=None, height=None):
    """Add an image to a slide at specified position."""
    if os.path.exists(image_path):
        kwargs = {"left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        slide.shapes.add_picture(image_path, **kwargs)
        return True
    return False


def build():
    prs = Presentation(TEMPLATE)

    # Remove all existing template slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get("r:id")
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # ================================================================
    # SLIDE 1 — Title
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "Title Slide Smoke"))
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Adaptive Skill and Stake\nin Forecast Markets", font_size=32, bold=True)
        elif ph.placeholder_format.idx == 1:
            set_text(ph, "Coupling Self-Financed Wagering with Online Skill Learning", font_size=18, color=SLATE)
        elif ph.placeholder_format.idx == 10:
            set_text(ph, "Anastasia Cattaneo  •  Imperial College London  •  Dyson School of Design Engineering", font_size=11, color=SLATE)

    # ================================================================
    # SLIDE 2 — Why Forecast Aggregation Matters
    # ================================================================
    add_bullet_slide(prs, "Title and Content Smoke",
        title="Why Forecast Aggregation Matters",
        subtitle="PROBLEM",
        bullets=[
            "Better predictions → better decisions (energy, logistics, finance, policy)",
            "Information is distributed across many sources with private data",
            "Combining forecasts reduces error — well-established in the literature",
            "Modern standard: full probabilistic forecasts, not point estimates",
            "Quality measured by strictly proper scoring rules (e.g. CRPS)",
            "A scoring rule is strictly proper ⟹ truthful reporting is optimal",
            "Open question: how to incentivise participation and weight forecasters?",
        ])

    # ================================================================
    # SLIDE 3 — Prediction Markets as a Solution
    # ================================================================
    add_bullet_slide(prs, "Title and Content",
        title="Prediction Markets as a Solution",
        subtitle="PROBLEM",
        bullets=[
            "Share predictions instead of raw data → reward based on quality",
            "Market aggregates individual forecasts into a collective view",
            "Real platforms: Numerai, Polymarket ($3.7B volume), Kalshi",
            "",
            "⚠ Polymarket: wash trading peaked at ~60% of weekly volume (Sirolly et al., 2025)",
            "⚠ Prices driven by small core of elite traders, not broad participation (Wu, 2025)",
            "",
            "Real markets are strategically adversarial → need formal guarantees",
        ])

    # ================================================================
    # SLIDE 4 — Existing Work (three columns)
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "Three Content"))
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Existing Work", font_size=28, bold=True)
        elif ph.placeholder_format.idx == 13:
            set_text(ph, "PROBLEM", font_size=14, color=SLATE)
        elif ph.placeholder_format.idx == 1:
            # Column 1: Self-Financed Wagering
            ph.text = ""
            tf = ph.text_frame
            tf.clear()
            lines = [
                ("Self-Financed Wagering", True, 14),
                ("Lambert et al. (2008)", False, 11),
                ("• Budget balance, truthfulness,", False, 10),
                ("  sybilproofness, ind. rationality", False, 10),
                ("• Unique mechanism (5 properties)", False, 10),
                ("Raja et al. (2024)", False, 11),
                ("• Added client + utility component", False, 10),
                ("• QA sharper than linear pooling", False, 10),
                ("", False, 10),
                ("⚠ History-free", True, 10),
            ]
            for i, (text, bold, size) in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = text
                p.font.size = Pt(size)
                p.font.bold = bold
                p.font.color.rgb = NAVY if bold else SLATE
        elif ph.placeholder_format.idx == 14:
            # Column 2: Online Aggregation
            ph.text = ""
            tf = ph.text_frame
            tf.clear()
            lines = [
                ("Online Forecast Aggregation", True, 14),
                ("Online convex optimisation", False, 11),
                ("• OGD, Hedge algorithms", False, 10),
                ("• Learns time-varying weights", False, 10),
                ("• Formal regret guarantees", False, 10),
                ("", False, 10),
                ("", False, 10),
                ("", False, 10),
                ("", False, 10),
                ("⚠ Non-strategic, always available", True, 10),
            ]
            for i, (text, bold, size) in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = text
                p.font.size = Pt(size)
                p.font.bold = bold
                p.font.color.rgb = NAVY if bold else SLATE
        elif ph.placeholder_format.idx == 15:
            # Column 3: Intermittent Contributions
            ph.text = ""
            tf = ph.text_frame
            tf.clear()
            lines = [
                ("Intermittent Contributions", True, 14),
                ("Vitali & Pinson (2025)", False, 11),
                ("• Robust regression + correction", False, 10),
                ("  matrix for missing forecasts", False, 10),
                ("• Shapley + out-of-sample payoff", False, 10),
                ("• Real wind energy case study", False, 10),
                ("  (Belgium, 9 NWP-based sellers)", False, 10),
                ("", False, 10),
                ("", False, 10),
                ("⚠ Relative weights on simplex", True, 10),
            ]
            for i, (text, bold, size) in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = text
                p.font.size = Pt(size)
                p.font.bold = bold
                p.font.color.rgb = NAVY if bold else SLATE

    # ================================================================
    # SLIDE 5 — The Gap and Contribution
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "Large Text Blue"))
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Stake alone is not a good proxy\nfor information quality", font_size=28, bold=True, color=WHITE)

    # ================================================================
    # SECTION DIVIDER — SOLUTION
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "Section Header Smoke"))
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Solution", font_size=36, bold=True)
        elif ph.placeholder_format.idx == 1:
            set_text(ph, "Mechanism Design", font_size=18, color=SLATE)

    # ================================================================
    # SLIDE 6 — Round-by-Round Mechanism
    # ================================================================
    add_bullet_slide(prs, "Title and Content Smoke",
        title="Round-by-Round Mechanism",
        subtitle="SOLUTION",
        bullets=[
            "① Submit: forecast (quantiles) + deposit from wealth × confidence",
            "     b_i = min(W_i,  f · W_i · c_i,  b_max)",
            "",
            "② Skill gate: effective wager = deposit × skill factor",
            "     m_i = b_i · (λ + (1−λ) · σ_i^η)     refund = b_i − m_i",
            "",
            "③ Aggregate: weighted combination using m_i as weights",
            "④ Settle: Π_i = m_i · (1 + s_i − s̄)     Σ Π_i = Σ m_i",
            "⑤ Update: EWMA loss → new σ_i;  W_{t+1} = max(0, W_t + profit)",
            "",
            "Same m_i for both influence and exposure → incentives aligned",
            "σ_i fixed BEFORE round t, from past losses only → truthfulness preserved",
        ])

    # ================================================================
    # SLIDE 7 — Skill Signal (with plot)
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "The Skill Signal", font_size=28, bold=True)
        elif ph.placeholder_format.idx == 13:
            set_text(ph, "SOLUTION", font_size=14, color=SLATE)
        elif ph.placeholder_format.idx == 1:
            ph.text = ""
            tf = ph.text_frame
            tf.clear()
            lines = [
                "Present:  L_t = (1−ρ)·L_{t−1} + ρ·loss_t",
                "Absent:   L_t = (1−κ)·L_{t−1} + κ·L₀",
                "Mapping:  σ = σ_min + (1−σ_min)·exp(−γ·L)",
                "",
                "Properties:",
                "• Absolute — reliability independent of other participants",
                "• Pre-round — computed from past losses only",
                "• Handles intermittency — absent forecasters decay to baseline",
                "",
                "Rank correlation (learned vs true): 1.0000 (T=20000, 20 seeds)",
            ]
            for i, text in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = text
                p.font.size = Pt(13)
                p.font.color.rgb = SLATE

    # ================================================================
    # SLIDE 8 — Architecture
    # ================================================================
    add_bullet_slide(prs, "Title and Content Smoke",
        title="Architecture and Implementation",
        subtitle="SOLUTION",
        bullets=[
            "Three-layer separation:",
            "  Layer 1 — Environment: DGPs (exogenous / endogenous truth)",
            "  Layer 2 — Agents: behaviour block (honest, noisy, sybil, arbitrageur, ...)",
            "  Layer 3 — Platform: core mechanism (scoring → aggregation → settlement → skill)",
            "",
            "Contract: agents output (participate, report, deposit). Core does not know motives.",
            "",
            "onlinev2 Python package — 20+ invariant tests, both point and quantile modes",
            "Experiment ladder: correctness → forecasting → dynamic → strategic robustness",
        ])

    # ================================================================
    # SECTION DIVIDER — VALIDATION
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "Section Header Blue"))
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Validation", font_size=36, bold=True, color=WHITE)
        elif ph.placeholder_format.idx == 1:
            set_text(ph, "Experimental Results", font_size=18, color=WHITE)

    # ================================================================
    # SLIDE 9 — Correctness
    # ================================================================
    add_bullet_slide(prs, "Title and Content",
        title="Correctness: The Mechanism Works",
        subtitle="VALIDATION — Rung 1",
        bullets=[
            "Budget balance:  max gap = 2.84 × 10⁻¹⁴  (1000 rounds)",
            "Zero-sum:  mean profit = 3.01 × 10⁻¹⁷",
            "Equal-score tie:  zero profit ✓",
            "Sybil (identical reports):  ratio = 1.000000,  max|Δ| = 2.07 × 10⁻¹⁷",
            "Noise–skill correlation:  −0.952 (MAE),  −0.979 (CRPS)",
            "",
            "Scoring tests: pinball ≥ 0 ✓,  CRPS ≥ 0 ✓,  perfect > shifted ✓,  CRPS ≤ 2 ✓",
            "All 20+ unit tests PASS (both modes)",
        ])

    # ================================================================
    # SLIDE 10 — Deposit Design (with plot)
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "Title, Text and Content"))
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Deposit Design Is the Strongest Lever", font_size=24, bold=True)
        elif ph.placeholder_format.idx == 13:
            set_text(ph, "VALIDATION — Rung 2", font_size=14, color=SLATE)
        elif ph.placeholder_format.idx == 16:
            ph.text = ""
            tf = ph.text_frame
            tf.clear()
            lines = [
                "IID Exponential:     0.0456 ± 0.0003",
                "Fixed Unit:            0.0423 ± 0.0002",
                "Bankroll+Conf:      0.0375 ± 0.0001  (−11.3%)",
                "Oracle Precision:   0.0227 ± 0.0001  (−46.3%)",
                "",
                "How stake enters the system matters",
                "more than the weighting rule.",
            ]
            for i, text in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = text
                p.font.size = Pt(12)
                p.font.color.rgb = SLATE
                if "matters" in text:
                    p.font.bold = True
                    p.font.color.rgb = NAVY
    # Add the plot
    plot_path = os.path.join(PLOTS, "deposit_policy_comparison.png")
    if os.path.exists(plot_path):
        # Place in the content placeholder area (right side)
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 17:
                # Insert image into the content placeholder
                ph.insert_picture(plot_path)
                break

    # ================================================================
    # SLIDE 11 — Weight Rules (with plot)
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "Title, Text and Content"))
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Weight Rules and the Combination Puzzle", font_size=24, bold=True)
        elif ph.placeholder_format.idx == 13:
            set_text(ph, "VALIDATION — Rung 2", font_size=14, color=SLATE)
        elif ph.placeholder_format.idx == 16:
            ph.text = ""
            tf = ph.text_frame
            tf.clear()
            lines = [
                "Fixed deposits (isolate skill):",
                "  Uniform:      0.0434",
                "  Skill only:    0.0419  (−3.5%)",
                "  Mechanism:  0.0423",
                "  Best single:  0.0232",
                "",
                "Bankroll deposits:",
                "  Deposit only: 0.0230",
                "  Mechanism:   0.0375",
                "",
                "Equal weights are hard to beat.",
                "Gains conditional on heterogeneity.",
            ]
            for i, text in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = text
                p.font.size = Pt(11)
                p.font.color.rgb = SLATE
                if "hard to beat" in text:
                    p.font.bold = True
                    p.font.color.rgb = NAVY
    plot_path = os.path.join(PLOTS, "weight_rule_comparison.png")
    if os.path.exists(plot_path):
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 17:
                ph.insert_picture(plot_path)
                break

    # ================================================================
    # SLIDE 12 — Skill Recovery (with plot)
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "Title, Text and Content"))
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Skill Recovery and Dynamic Robustness", font_size=24, bold=True)
        elif ph.placeholder_format.idx == 13:
            set_text(ph, "VALIDATION — Rung 3", font_size=14, color=SLATE)
        elif ph.placeholder_format.idx == 16:
            ph.text = ""
            tf = ph.text_frame
            tf.clear()
            lines = [
                "6 forecasters, T=20000, 20 seeds:",
                "  f0 (τ=0.15): σ = 0.959",
                "  f1 (τ=0.22): σ = 0.942",
                "  f2 (τ=0.32): σ = 0.919",
                "  f3 (τ=0.46): σ = 0.890",
                "  f4 (τ=0.68): σ = 0.854",
                "  f5 (τ=1.00): σ = 0.820",
                "",
                "Spearman ρ(τ, σ) = 1.0000",
                "Perfect rank recovery.",
                "",
                "Staleness decay (κ > 0) punishes",
                "strategic absence.",
            ]
            for i, text in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = text
                p.font.size = Pt(11)
                p.font.color.rgb = SLATE
                if "1.0000" in text or "Perfect" in text:
                    p.font.bold = True
                    p.font.color.rgb = NAVY
    plot_path = os.path.join(PLOTS, "quantiles_crps_recovery.png")
    if os.path.exists(plot_path):
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 17:
                ph.insert_picture(plot_path)
                break

    # ================================================================
    # SLIDE 13 — Strategic Robustness (with plot)
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "Title, Text and Content"))
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Strategic Robustness", font_size=28, bold=True)
        elif ph.placeholder_format.idx == 13:
            set_text(ph, "VALIDATION — Rung 4", font_size=14, color=SLATE)
        elif ph.placeholder_format.idx == 16:
            ph.text = ""
            tf = ph.text_frame
            tf.clear()
            lines = [
                "Sybil (identical reports):",
                "  Profit ratio = 1.000000",
                "  max|Δ| = 2.07 × 10⁻¹⁷",
                "  → No advantage from splitting",
                "",
                "Sybil (diversified reports):",
                "  Ratio = 1.065",
                "  → Sybilproofness requires identical reports",
                "",
                "Arbitrage scan (all λ values):",
                "  Profit = 0.0 across all configurations",
                "",
                "Adaptive adversaries remain open.",
            ]
            for i, text in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = text
                p.font.size = Pt(11)
                p.font.color.rgb = SLATE
                if "No advantage" in text or "Profit = 0.0" in text:
                    p.font.bold = True
                    p.font.color.rgb = NAVY
    plot_path = os.path.join(PLOTS, "sybil.png")
    if os.path.exists(plot_path):
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 17:
                ph.insert_picture(plot_path)
                break

    # ================================================================
    # SLIDE 14 — Contributions, Limitations, Closing
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "Two Content"))
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Contributions and Limitations", font_size=28, bold=True)
        elif ph.placeholder_format.idx == 13:
            set_text(ph, "", font_size=14, color=SLATE)
        elif ph.placeholder_format.idx == 1:
            # Left column: Contributions
            ph.text = ""
            tf = ph.text_frame
            tf.clear()
            lines = [
                ("Contributions", True, 16),
                ("1. Wagering + online skill learning", False, 11),
                ("    Absolute, pre-round, handles intermittency", False, 10),
                ("2. Verified: budget balance < 10⁻¹⁴,", False, 11),
                ("    sybilproofness ratio = 1.000000", False, 10),
                ("3. Deposit design strongest lever:", False, 11),
                ("    bankroll-confidence −11.3% vs fixed", False, 10),
                ("4. Skill recovery: ρ = 1.0000", False, 11),
                ("5. Sybil-resistant, no arbitrage profit", False, 11),
                ("6. onlinev2 platform + dashboard", False, 11),
            ]
            for i, (text, bold, size) in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = text
                p.font.size = Pt(size)
                p.font.bold = bold
                p.font.color.rgb = NAVY if bold else SLATE
        elif ph.placeholder_format.idx == 2:
            # Right column: Limitations
            ph.text = ""
            tf = ph.text_frame
            tf.clear()
            lines = [
                ("Limitations", True, 16),
                ("• Tail calibration: under-dispersion", False, 11),
                ("  ~5 pp in tails (quantile averaging)", False, 10),
                ("• Equal weights competitive —", False, 11),
                ("  gains conditional on heterogeneity", False, 10),
                ("• Truthfulness under risk neutrality", False, 11),
                ("  only (Lambert assumption)", False, 10),
                ("• All synthetic data — no real-world", False, 11),
                ("  deployment yet", False, 10),
                ("", False, 10),
                ("Strongest lever: deposit design.", True, 12),
            ]
            for i, (text, bold, size) in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = text
                p.font.size = Pt(size)
                p.font.bold = bold
                p.font.color.rgb = NAVY if bold else SLATE

    # ================================================================
    # SLIDE 15 — Thank You
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "Closing Slide Blue"))
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Thank you", font_size=36, bold=True, color=WHITE)
        elif ph.placeholder_format.idx == 10:
            set_text(ph, "Anastasia Cattaneo  •  Imperial College London  •  2025", font_size=14, color=WHITE)

    # ================================================================
    # Save
    # ================================================================
    prs.save(OUTPUT)
    print(f"Saved {OUTPUT} with {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
