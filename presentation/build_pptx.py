"""
Build thesis defence PowerPoint from Imperial College template.
Academic formatting: clear hierarchy, bold headings, proper fonts.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import copy
from lxml import etree

# Paths
TEMPLATE = "presentation/template_original.pptx"
OUTPUT = "presentation/thesis_defence.pptx"
PLOTS_DIR = "onlinev2/outputs_final/core/experiments"
PLOTS_DIR_OLD = "onlinev2/outputs/core/experiments"

# Layout indices
LAYOUT_TITLE_SMOKE = 1
LAYOUT_SECTION_BLUE = 28
LAYOUT_CONTENT = 9
LAYOUT_CONTENT_SMOKE = 10
LAYOUT_TWO_CONTENT = 12
LAYOUT_THREE_CONTENT = 15
LAYOUT_LARGE_TEXT_BLUE = 24
LAYOUT_TITLE_TEXT_CONTENT = 29
LAYOUT_CLOSING_BLUE = 56
LAYOUT_TITLE_ONLY = 57
LAYOUT_BLANK = 60

# Colours (Imperial palette)
IMPERIAL_BLUE = RGBColor(0x00, 0x3E, 0x74)
IMPERIAL_DARK = RGBColor(0x00, 0x2A, 0x4E)
ACCENT_GREEN = RGBColor(0x02, 0x89, 0x3B)
ACCENT_RED = RGBColor(0xDD, 0x22, 0x20)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
MED_GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GREY = RGBColor(0x99, 0x99, 0x99)


def get_plot(name):
    """Find a plot file, checking outputs_final first, then outputs."""
    path = os.path.join(PLOTS_DIR, name)
    if os.path.exists(path):
        return path
    path = os.path.join(PLOTS_DIR_OLD, name)
    if os.path.exists(path):
        return path
    return None


def set_text_bold(paragraph, text, font_size=Pt(14), color=DARK_GREY):
    """Set paragraph text with bold formatting."""
    paragraph.text = ""
    run = paragraph.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = font_size
    run.font.color.rgb = color


def add_formatted_bullets(tf, items, base_size=Pt(14)):
    """Add formatted bullet items. Items starting with '##' are bold subheadings."""
    tf.clear()
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        
        if item.startswith("##"):
            # Bold subheading
            text = item[2:].strip()
            p.text = ""
            run = p.add_run()
            run.text = text
            run.font.bold = True
            run.font.size = base_size
            run.font.color.rgb = IMPERIAL_BLUE
            p.level = 0
            p.space_before = Pt(8)
        elif item.startswith(">>"):
            # Highlighted key result
            text = item[2:].strip()
            p.text = ""
            run = p.add_run()
            run.text = text
            run.font.bold = True
            run.font.size = Pt(base_size.pt + 1)
            run.font.color.rgb = ACCENT_GREEN
            p.level = 0
        elif item.startswith("!!"):
            # Warning/limitation
            text = item[2:].strip()
            p.text = ""
            run = p.add_run()
            run.text = text
            run.font.bold = False
            run.font.size = Pt(base_size.pt - 1)
            run.font.color.rgb = ACCENT_RED
            p.level = 1
        elif item == "":
            p.text = ""
            p.space_before = Pt(4)
        elif item.startswith("  "):
            # Indented sub-bullet
            p.text = item.strip()
            p.level = 2
            p.font.size = Pt(base_size.pt - 2)
            p.font.color.rgb = MED_GREY
        else:
            p.text = item
            p.level = 1
            p.font.size = base_size
            p.font.color.rgb = DARK_GREY


def add_image_to_slide(slide, image_path, left, top, width=None, height=None):
    """Add image to slide at specified position."""
    if width and height:
        slide.shapes.add_picture(image_path, left, top, width, height)
    elif width:
        slide.shapes.add_picture(image_path, left, top, width=width)
    elif height:
        slide.shapes.add_picture(image_path, left, top, height=height)
    else:
        slide.shapes.add_picture(image_path, left, top)


def build_presentation():
    prs = Presentation(TEMPLATE)
    
    # Remove all existing template slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # =========================================================================
    # SLIDE 1 — Title
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_SMOKE])
    slide.placeholders[0].text = "Adaptive Skill and Stake\nin Forecast Markets"
    slide.placeholders[1].text = "Coupling Self-Financed Wagering with Online Skill Learning"
    slide.placeholders[10].text = "Anastasia Cattaneo\nImperial College London — Dyson School of Design Engineering\n2025"

    # =========================================================================
    # SLIDE 2 — Why Forecast Aggregation Matters
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    slide.placeholders[0].text = "Why Forecast Aggregation Matters"
    slide.placeholders[13].text = ""
    tf = slide.placeholders[1].text_frame
    add_formatted_bullets(tf, [
        "##Combining forecasts reduces error",
        "Different sources capture different aspects of reality",
        "Averaging out individual errors improves collective accuracy",
        "",
        "##Modern standard: probabilistic forecasts",
        "Full distributions, not point estimates (Gneiting & Raftery, 2007)",
        "Quality measured by strictly proper scoring rules (e.g. CRPS)",
        "Strictly proper → only way to maximise score is to report true belief",
        "",
        "##Open question",
        "How to incentivise participation?",
        "How to decide whose forecast should count more?",
    ], base_size=Pt(14))

    # =========================================================================
    # SLIDE 3 — Prediction Markets
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    slide.placeholders[0].text = "Prediction Markets as a Solution"
    slide.placeholders[13].text = ""
    tf = slide.placeholders[1].text_frame
    add_formatted_bullets(tf, [
        "##Mechanism",
        "Data owners share predictions (not raw data)",
        "Rewarded based on forecast accuracy after outcome observed",
        "Market aggregates individual predictions into collective view",
        "",
        "##Real platforms",
        "Numerai (data science tournaments), Polymarket, Kalshi",
        "",
        "##Evidence of problems",
        "!!Wash trading ~60% of Polymarket weekly volume (Sirolly et al., 2025)",
        "!!Prices driven by small core of active traders, not broad wisdom (Wu, 2025)",
        "",
        ">>→ Need mechanisms with formal guarantees",
    ], base_size=Pt(14))

    # =========================================================================
    # SLIDE 4 — Existing Work (Three columns)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_THREE_CONTENT])
    slide.placeholders[0].text = "Existing Work"
    slide.placeholders[13].text = ""
    
    # Column 1: Self-Financed Wagering
    tf1 = slide.placeholders[1].text_frame
    add_formatted_bullets(tf1, [
        "##Self-Financed Wagering",
        "Lambert et al. (2008)",
        "Weighted-score mechanism",
        "7 properties; uniqueness result",
        "",
        "Raja et al. (2024)",
        "Added client + utility",
        "Conditionally truthful",
        "Wind energy case study",
        "",
        "!!Limitation: history-free",
    ], base_size=Pt(11))

    # Column 2: Online Aggregation
    tf2 = slide.placeholders[14].text_frame
    add_formatted_bullets(tf2, [
        "##Online Aggregation",
        "OGD / Hedge algorithms",
        "Learns time-varying weights",
        "Regret guarantees",
        "",
        "",
        "",
        "",
        "",
        "!!Non-strategic assumption",
        "!!No payments or incentives",
    ], base_size=Pt(11))

    # Column 3: Intermittent Contributions
    tf3 = slide.placeholders[15].text_frame
    add_formatted_bullets(tf3, [
        "##Intermittent Contributions",
        "Vitali & Pinson (2025)",
        "Correction matrix for missing",
        "OGD on pinball loss",
        "Shapley + scoring payoff",
        "Belgian offshore wind",
        "",
        "",
        "",
        "!!Relative weights on simplex",
        "!!Different settlement structure",
    ], base_size=Pt(11))

    # =========================================================================
    # SLIDE 5 — Gap and Contribution
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_LARGE_TEXT_BLUE])
    # Large text layout only has title placeholder
    title_tf = slide.placeholders[0].text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    p.text = ""
    run = p.add_run()
    run.text = "No existing design couples self-financed wagering with an online skill-learning layer"
    run.font.size = Pt(24)
    run.font.bold = True
    
    p2 = title_tf.add_paragraph()
    p2.text = ""
    p2.space_before = Pt(20)
    run2 = p2.add_run()
    run2.text = "Contribution: effective wager = deposit × learned skill"
    run2.font.size = Pt(20)
    run2.font.bold = True
    
    p3 = title_tf.add_paragraph()
    p3.space_before = Pt(16)
    p3.text = ""
    run3 = p3.add_run()
    run3.text = "• Absolute (not relative)  • Pre-round (preserves truthfulness)  • Handles intermittency"
    run3.font.size = Pt(16)

    # =========================================================================
    # SECTION DIVIDER — SOLUTION
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION_BLUE])
    slide.placeholders[0].text = "Solution"
    slide.placeholders[1].text = "Mechanism Design and Implementation"

    # =========================================================================
    # SLIDE 6 — Round-by-Round Mechanism
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    slide.placeholders[0].text = "The Mechanism: Round-by-Round"
    slide.placeholders[13].text = ""
    tf = slide.placeholders[1].text_frame
    add_formatted_bullets(tf, [
        "##Five steps per round:",
        "",
        "##1. Submit",
        "Player submits quantile forecast + deposit (wealth × confidence)",
        "",
        "##2. Skill Gate",
        "Effective wager = deposit × skill factor; remainder refunded",
        "",
        "##3. Aggregate",
        "Weighted combination using effective wagers as weights",
        "",
        "##4. Settle",
        "Payoff redistributes wager pool by relative scores (budget-balanced)",
        "",
        "##5. Update",
        "Loss → EWMA → skill recomputed; wealth updated",
        "",
        ">>Same effective wager controls BOTH influence and exposure",
    ], base_size=Pt(13))

    # =========================================================================
    # SLIDE 7 — Skill Signal (text + image)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_TEXT_CONTENT])
    slide.placeholders[0].text = "The Skill Signal"
    slide.placeholders[13].text = "Core Innovation"
    
    tf = slide.placeholders[16].text_frame
    add_formatted_bullets(tf, [
        "##When present:",
        "Loss state blends previous",
        "value with current round loss",
        "(EWMA with learning rate ρ)",
        "",
        "##When absent:",
        "Loss state reverts toward",
        "neutral baseline (staleness decay)",
        "",
        "##Mapping:",
        "Exponential: loss → skill ∈ [σ_min, 1]",
        "",
        "##Key difference from Vitali-Pinson:",
        "Absolute skill (not relative)",
        "One forecaster improves without",
        "reducing another's weight",
    ], base_size=Pt(11))
    
    # Add skill recovery plot on right side
    plot_path = get_plot("skill_recovery/plots/quantiles_crps_recovery.png")
    if plot_path:
        # Tall plot (2096x1491), constrain by height
        add_image_to_slide(slide, plot_path,
                          left=Inches(5.0), top=Inches(1.3),
                          height=Inches(5.0))

    # =========================================================================
    # SLIDE 8 — Architecture
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    slide.placeholders[0].text = "Architecture and Implementation"
    slide.placeholders[13].text = ""
    tf = slide.placeholders[1].text_frame
    add_formatted_bullets(tf, [
        "##Three-layer modular separation:",
        "",
        "##Environment",
        "Data-generating processes (exogenous/endogenous truth)",
        "",
        "##Agents",
        "Behaviour policies: honest, noisy, risk-averse, sybil, arbitrageur, colluder, wash trader",
        "Each outputs: (participate?, report, deposit) — core does not know motives",
        "",
        "##Platform (deterministic, side-effect-free)",
        "Scoring → Aggregation → Settlement → Skill update",
        "",
        "##Implementation: onlinev2 Python package",
        "20+ invariant tests; property-based testing (Hypothesis)",
        "Experiment ladder: correctness → forecasting → dynamics → strategy",
    ], base_size=Pt(13))

    # =========================================================================
    # SECTION DIVIDER — VALIDATION
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION_BLUE])
    slide.placeholders[0].text = "Validation"
    slide.placeholders[1].text = "Experimental Results"

    # =========================================================================
    # SLIDE 9 — Correctness (text + settlement plot)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_TEXT_CONTENT])
    slide.placeholders[0].text = "Correctness: The Mechanism Works"
    slide.placeholders[13].text = "1000 rounds, 20 seeds"
    
    tf = slide.placeholders[16].text_frame
    add_formatted_bullets(tf, [
        "##Budget Balance",
        "Max gap: 2.84 × 10⁻¹⁴",
        "Mean profit: 3.01 × 10⁻¹⁷",
        "  (machine precision)",
        "",
        "##Sybilproofness",
        "Identical reports ratio: 1.000000",
        "Max |Δ|: 2.07 × 10⁻¹⁷",
        "",
        "##Scoring Invariants",
        "Pinball ≥ 0 ✓",
        "CRPS ≥ 0 ✓",
        "Perfect beats shifted ✓",
        "CRPS bounded ✓",
        "",
        ">>All 20+ tests PASS",
    ], base_size=Pt(11))
    
    # Add settlement sanity plot
    plot_path = get_plot("settlement_sanity/plots/settlement_sanity.png")
    if plot_path:
        # Tall plot (1945x1343), constrain by height
        add_image_to_slide(slide, plot_path,
                          left=Inches(5.0), top=Inches(1.3),
                          height=Inches(5.0))

    # =========================================================================
    # SLIDE 10 — Deposit Design (KEY RESULT)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_TEXT_CONTENT])
    slide.placeholders[0].text = "Deposit Design Is the Strongest Lever"
    slide.placeholders[13].text = "Key Finding"
    
    tf = slide.placeholders[16].text_frame
    add_formatted_bullets(tf, [
        "##Deposit policy comparison",
        "  (mechanism weight rule, 20 seeds)",
        "",
        "IID Exponential:    0.0456 ± 0.0003",
        "Fixed Unit (b=1):   0.0423 ± 0.0002",
        "Bankroll + Conf:    0.0375 ± 0.0001",
        "Oracle Precision:   0.0227 ± 0.0001",
        "",
        ">>Bankroll vs Fixed: −11.3%",
        ">>Oracle vs Fixed:   −46.3%",
        "",
        "##Implication:",
        "How stake enters the system",
        "matters more than the",
        "weighting rule",
    ], base_size=Pt(12))
    
    plot_path = get_plot("deposit_policy_comparison/plots/deposit_policy_comparison.png")
    if plot_path:
        # Wide plot (2095x903, ratio 2.3:1), use full right area
        add_image_to_slide(slide, plot_path,
                          left=Inches(4.6), top=Inches(2.2),
                          width=Inches(7.5))

    # =========================================================================
    # SLIDE 11 — Weight Rules
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_TEXT_CONTENT])
    slide.placeholders[0].text = "Weight Rules and the Combination Puzzle"
    slide.placeholders[13].text = ""
    
    tf = slide.placeholders[16].text_frame
    add_formatted_bullets(tf, [
        "##Under fixed deposits:",
        "Uniform:     0.0434 ± 0.0002",
        "Deposit:     0.0434 ± 0.0002",
        "Skill:       0.0419 ± 0.0002",
        "Mechanism:   0.0423 ± 0.0002",
        "Best single: 0.0232 ± 0.0001",
        "",
        ">>Skill vs uniform: −3.5%",
        "",
        "##Under bankroll deposits:",
        "Deposit only: 0.0230 ± 0.0001",
        "Mechanism:    0.0375 ± 0.0001",
        "",
        "##Forecast combination puzzle:",
        "Equal weights hard to beat",
        "(Wang et al., 2023)",
    ], base_size=Pt(11))
    
    plot_path = get_plot("weight_rule_comparison/plots/weight_rule_comparison.png")
    if plot_path:
        # Tall plot (2094x1491), constrain by height
        add_image_to_slide(slide, plot_path,
                          left=Inches(5.0), top=Inches(1.3),
                          height=Inches(5.0))

    # =========================================================================
    # SLIDE 12 — Skill Recovery
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    slide.placeholders[0].text = "Skill Recovery and Dynamic Robustness"
    slide.placeholders[13].text = "T=20000, 6 forecasters, 20 seeds"
    tf = slide.placeholders[1].text_frame
    add_formatted_bullets(tf, [
        "##Skill recovery results (quantiles_crps mode):",
        "",
        "  f0: τ=0.15  →  loss=0.023  →  σ=0.959",
        "  f1: τ=0.22  →  loss=0.033  →  σ=0.942",
        "  f2: τ=0.32  →  loss=0.047  →  σ=0.919",
        "  f3: τ=0.46  →  loss=0.066  →  σ=0.890",
        "  f4: τ=0.68  →  loss=0.089  →  σ=0.854",
        "  f5: τ=1.00  →  loss=0.112  →  σ=0.820",
        "",
        ">>Spearman rank correlation (τ vs σ): 1.0000",
        ">>Perfect rank recovery in both point and quantile modes",
        "",
        "##Staleness decay:",
        "Absent forecasters revert toward baseline",
        "→ removes incentive for strategic absence",
    ], base_size=Pt(13))

    # =========================================================================
    # SLIDE 13 — Strategic Robustness
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_TEXT_CONTENT])
    slide.placeholders[0].text = "Strategic Robustness"
    slide.placeholders[13].text = ""
    
    tf = slide.placeholders[16].text_frame
    add_formatted_bullets(tf, [
        "##Sybil attacks (identical reports):",
        "Profit ratio: 1.000000",
        "Max |Δ|: 2.07 × 10⁻¹⁷",
        ">>No advantage from splitting",
        "",
        "##Sybil (diversified reports):",
        "Ratio: 1.065",
        "!!Slight advantage (known limit)",
        "",
        "##Strategic deposit manipulation:",
        "Ratio: 1.000000",
        "",
        "##Arbitrage (Chen et al., 2014):",
        "Zero profit in repeated setting",
        "Skill gate limits sustained exploit",
        "",
        "!!Adaptive adversaries: open challenge",
    ], base_size=Pt(11))
    
    plot_path = get_plot("sybil/plots/sybil.png")
    if plot_path:
        # Very wide plot (2695x829, ratio 3.25:1), use full right area
        add_image_to_slide(slide, plot_path,
                          left=Inches(4.6), top=Inches(2.5),
                          width=Inches(7.8))

    # =========================================================================
    # SLIDE 14 — Calibration
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_TEXT_CONTENT])
    slide.placeholders[0].text = "Calibration"
    slide.placeholders[13].text = "Known Limitation"
    
    tf = slide.placeholders[16].text_frame
    add_formatted_bullets(tf, [
        "##Reliability (latent-fixed DGP)",
        "  T=20000, 10 forecasters, 20 seeds",
        "",
        "τ=0.10:  p̂=0.054  (dev: −0.046)",
        "τ=0.25:  p̂=0.194  (dev: −0.056)",
        ">>τ=0.50:  p̂=0.499  (dev: −0.001)",
        "τ=0.75:  p̂=0.804  (dev: +0.054)",
        "τ=0.90:  p̂=0.945  (dev: +0.045)",
        "",
        "##Median: nearly perfect",
        "!!Tails: ~5pp systematic under-dispersion",
        "",
        "##Inherent to quantile averaging",
        "Shared across all weighting methods",
        "Not specific to this mechanism",
    ], base_size=Pt(11))
    
    plot_path = get_plot("calibration/plots/calibration_reliability.png")
    if plot_path:
        # Nearly square plot (913x895), constrain by height
        add_image_to_slide(slide, plot_path,
                          left=Inches(5.5), top=Inches(1.3),
                          height=Inches(5.0))

    # =========================================================================
    # SLIDE 15 — Contributions and Limitations
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TWO_CONTENT])
    slide.placeholders[0].text = "Contributions and Limitations"
    slide.placeholders[13].text = ""
    
    # Left column: Contributions
    tf1 = slide.placeholders[1].text_frame
    add_formatted_bullets(tf1, [
        "##Contributions",
        "",
        "1. Mechanism coupling wagering",
        "   with online skill learning",
        "2. Budget balance < 10⁻¹⁴",
        "   Sybilproof (ratio = 1.000000)",
        "3. Deposit design: −11.3% CRPS",
        "4. Skill recovery: ρ = 1.0000",
        "5. Sybil-resistant; no arbitrage",
        "6. Modular platform + test suite",
    ], base_size=Pt(12))

    # Right column: Limitations
    tf2 = slide.placeholders[2].text_frame
    add_formatted_bullets(tf2, [
        "##Limitations",
        "",
        "!!Tail calibration: ~5pp under-dispersion",
        "!!Equal weights competitive in some configs",
        "!!Truthfulness: risk-neutrality assumption",
        "!!All synthetic data — no deployment",
        "",
        "",
        "##Key takeaway:",
        ">>Strongest lever is deposit design,",
        ">>not the weighting rule",
    ], base_size=Pt(12))

    # =========================================================================
    # CLOSING SLIDE
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CLOSING_BLUE])
    slide.placeholders[0].text = "Thank you"
    slide.placeholders[10].text = (
        "Anastasia Cattaneo\n"
        "Imperial College London — Dyson School of Design Engineering\n"
        "2025"
    )

    # =========================================================================
    # Save
    # =========================================================================
    prs.save(OUTPUT)
    print(f"✓ Saved: {OUTPUT}")
    print(f"  Total slides: {len(prs.slides)}")
    
    # Print summary
    for i, slide in enumerate(prs.slides, 1):
        layout = slide.slide_layout.name
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text[:50].replace("\n", " ")
                if text.strip():
                    title = text.strip()
                    break
        imgs = sum(1 for s in slide.shapes if s.shape_type == 13)
        img_note = f" [{imgs} plot(s)]" if imgs else ""
        print(f"  {i:2d}. {title}{img_note}")


if __name__ == "__main__":
    build_presentation()
